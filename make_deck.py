"""
Generates recognize-deck.pptx — 6-slide pitch deck for Recognize AI.
Run: python3 make_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette (matches the app's design system) ──────────────────────────
BG       = RGBColor(0x06, 0x08, 0x10)   # deep space
BG2      = RGBColor(0x0C, 0x11, 0x20)   # slightly lighter dark
CLAY     = RGBColor(0xB8, 0x42, 0x2E)   # Boston Clay — primary accent
CYAN     = RGBColor(0x00, 0xD4, 0xFF)   # data / tech nodes
PURPLE   = RGBColor(0x9B, 0x6D, 0xFF)   # knowledge graph nodes
GOLD     = RGBColor(0xC8, 0xB8, 0x40)   # community detection
TEXT     = RGBColor(0xF7, 0xF5, 0xF2)   # warm white
MUTED    = RGBColor(0x6C, 0x72, 0x78)   # slate grey

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper: solid background ───────────────────────────────────────────────────
def fill_bg(slide, colour=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


# ── Helper: add a text box ─────────────────────────────────────────────────────
def add_text(slide, text, x, y, w, h,
             size=18, bold=False, colour=TEXT, align=PP_ALIGN.LEFT,
             font="Calibri", italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = colour
    run.font.name    = font
    return txb


# ── Helper: coloured rectangle ─────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, colour, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        x, y, w, h
    )
    shape.line.fill.background()   # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    return shape


# ── Helper: label + body text block ───────────────────────────────────────────
def add_label_body(slide, label, label_colour, title, body,
                   x, y, w, h_label=Inches(0.28), h_title=Inches(0.45), h_body=Inches(0.9)):
    add_text(slide, label, x, y,              w, h_label, size=9,  bold=True, colour=label_colour, font="Calibri")
    add_text(slide, title, x, y+h_label,      w, h_title, size=13, bold=True, colour=TEXT,         font="Calibri")
    add_text(slide, body,  x, y+h_label+h_title, w, h_body,  size=10, colour=MUTED,        font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
fill_bg(s1, BG)

# Left red accent bar
add_rect(s1, Inches(0), Inches(0), Inches(0.06), H, CLAY)

# Eyebrow
add_text(s1, "RECOGNIZE AI", Inches(0.5), Inches(1.2), Inches(8),
         Inches(0.35), size=11, bold=True, colour=CLAY, font="Calibri")

# Main title — 3 lines
add_text(s1, "The", Inches(0.5), Inches(1.65), Inches(9), Inches(1.0),
         size=72, bold=True, colour=TEXT, font="Calibri")
add_text(s1, "Company", Inches(0.5), Inches(2.5), Inches(9), Inches(1.1),
         size=72, bold=True, colour=CLAY, font="Calibri")
add_text(s1, "Brain.", Inches(0.5), Inches(3.4), Inches(9), Inches(1.0),
         size=72, bold=True, colour=TEXT, font="Calibri")

# Subtitle
add_text(s1,
    "Real-time meeting intelligence that identifies every speaker, "
    "transcribes every word, and builds a knowledge graph that never forgets.",
    Inches(0.5), Inches(4.6), Inches(7.5), Inches(0.9),
    size=15, colour=MUTED, font="Calibri")

# Pill tag
pill = add_rect(s1, Inches(0.5), Inches(5.7), Inches(3.8), Inches(0.42), BG2)
pill.line.color.rgb = CYAN
pill.line.width = Pt(0.75)
add_text(s1, "● Live diarisation  ·  GraphRAG  ·  Neo4j",
         Inches(0.62), Inches(5.72), Inches(3.6), Inches(0.38),
         size=10, colour=CYAN, font="Calibri")

# Slide number
add_text(s1, "01 / 06", Inches(12.1), Inches(7.0), Inches(1.1), Inches(0.35),
         size=9, colour=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
fill_bg(s2, BG2)
add_rect(s2, Inches(0), Inches(0), Inches(0.06), H, CLAY)

add_text(s2, "THE PROBLEM", Inches(0.5), Inches(0.5), Inches(8),
         Inches(0.3), size=10, bold=True, colour=CLAY)
add_text(s2, "A decades-old unsolved problem — solved.",
         Inches(0.5), Inches(0.9), Inches(12), Inches(0.7),
         size=32, bold=True, colour=TEXT)

rows = [
    ("01", "Speaker diarisation requires dedicated hardware",
     "Classical approaches need microphone arrays, per-speaker voice enrollment, or expensive "
     "cloud APIs with no real-time output. Single-camera, single-mic setups had no clean prior solution."),
    ("02", "Institutional knowledge evaporates after every meeting",
     "Decisions get made, context is lost, the same debate happens months later. "
     "Companies have no persistent, queryable memory of what they decided or why."),
    ("03", "Transcription, attribution, and memory have never been unified",
     "Entity extraction, cross-meeting contradiction detection, and graph-based RAG have existed "
     "separately in research — never in a single automatic pipeline, until now."),
]

y = Inches(1.75)
for num, title, body in rows:
    add_text(s2, num, Inches(0.5), y, Inches(0.4), Inches(0.35),
             size=11, bold=True, colour=CLAY)
    add_text(s2, title, Inches(1.0), y, Inches(11.5), Inches(0.35),
             size=13, bold=True, colour=TEXT)
    add_text(s2, body, Inches(1.0), y + Inches(0.37), Inches(11.5), Inches(0.55),
             size=11, colour=MUTED)
    y += Inches(1.1)

# Callout box
add_rect(s2, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), RGBColor(0x1C, 0x0A, 0x08))
add_rect(s2, Inches(0.5), Inches(6.0), Inches(0.07), Inches(0.9), CLAY)
add_text(s2,
    "Recognize sidesteps 30 years of hardware-dependent diarisation by watching lips "
    "instead of parsing voices — zero enrollment, works on the very first meeting.",
    Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.8),
    size=11, colour=TEXT)

add_text(s2, "02 / 06", Inches(12.1), Inches(7.0), Inches(1.1), Inches(0.35),
         size=9, colour=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
fill_bg(s3, BG)
add_rect(s3, Inches(0), Inches(0), Inches(0.06), H, CLAY)

add_text(s3, "HOW IT WORKS", Inches(0.5), Inches(0.5), Inches(8),
         Inches(0.3), size=10, bold=True, colour=CLAY)
add_text(s3, "Zero enrollment. Zero hardware. Works on the first meeting.",
         Inches(0.5), Inches(0.9), Inches(12.5), Inches(0.65),
         size=28, bold=True, colour=TEXT)

steps = [
    (CLAY,   "CAPTURE",  "Chrome\nExtension",  "Tab video + audio\nover WebSocket"),
    (CYAN,   "VISION",   "MediaPipe\nMAR",      "Lip movement per face\n→ active speaker"),
    (CLAY,   "ASR",      "Groq\nWhisper",       "Word-level timestamps\n6-second windows"),
    (PURPLE, "SYNC",     "Diarizer",            "Lip timeline ×\nword timestamps"),
    (CYAN,   "MEMORY",   "Neo4j\nGraph",        "Auto-ingested\non session end"),
]

box_w = Inches(2.35)
box_h = Inches(3.8)
gap   = Inches(0.18)
x0    = Inches(0.35)
y0    = Inches(1.7)

for i, (col, tag, name, desc) in enumerate(steps):
    bx = x0 + i * (box_w + gap)

    # Card background
    card = add_rect(s3, bx, y0, box_w, box_h, RGBColor(0x0C, 0x11, 0x20))
    card.line.color.rgb = RGBColor(0x1A, 0x22, 0x38)
    card.line.width = Pt(0.5)

    # Top colour bar
    add_rect(s3, bx, y0, box_w, Inches(0.06), col)

    # Colour dot
    dot = add_rect(s3, bx + Inches(0.2), y0 + Inches(0.22), Inches(0.28), Inches(0.28), col)

    # Tag
    add_text(s3, tag, bx + Inches(0.2), y0 + Inches(0.6), box_w - Inches(0.4), Inches(0.28),
             size=8, bold=True, colour=col)
    # Name
    add_text(s3, name, bx + Inches(0.2), y0 + Inches(0.95), box_w - Inches(0.4), Inches(0.75),
             size=15, bold=True, colour=TEXT)
    # Desc
    add_text(s3, desc, bx + Inches(0.2), y0 + Inches(1.8), box_w - Inches(0.4), Inches(1.6),
             size=11, colour=MUTED)

    # Arrow (except last)
    if i < len(steps) - 1:
        ax = bx + box_w + Inches(0.01)
        add_text(s3, "›", ax, y0 + Inches(1.6), Inches(0.18), Inches(0.5),
                 size=22, bold=True, colour=CLAY, align=PP_ALIGN.CENTER)

add_text(s3, "03 / 06", Inches(12.1), Inches(7.0), Inches(1.1), Inches(0.35),
         size=9, colour=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — KNOWLEDGE GRAPH / VECTOR SEARCH
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
fill_bg(s4, BG2)
add_rect(s4, Inches(0), Inches(0), Inches(0.06), H, CLAY)

add_text(s4, "THE KNOWLEDGE GRAPH", Inches(0.5), Inches(0.5), Inches(10),
         Inches(0.3), size=10, bold=True, colour=CLAY)
add_text(s4, "Vector memory that compounds across every meeting.",
         Inches(0.5), Inches(0.9), Inches(12.5), Inches(0.65),
         size=28, bold=True, colour=TEXT)

cards = [
    (CLAY,   "EMBEDDINGS",        "384-dim semantic vectors",
     "Every entity and chunk embedded with all-MiniLM-L6-v2. "
     "Cosine vector search finds related concepts even when different words are used across meetings."),
    (CYAN,   "DEDUPLICATION",     "Entities merge, never multiply",
     "New entities with ≥ 0.90 cosine similarity to existing ones are merged. "
     "Mention counts accumulate — the more something is discussed, the more central it becomes."),
    (PURPLE, "GRAPHRAG — LOCAL",  "Vector search + graph traversal",
     "Query → nearest entity nodes → 1-hop neighbor walk → source chunks → "
     "Claude returns a cited answer grounded in your actual meeting history."),
    (GOLD,   "GRAPHRAG — GLOBAL", "Louvain maps your knowledge terrain",
     "Community detection clusters tightly connected entities. Claude summarises each cluster. "
     "Holistic questions reason across the entire organisational graph."),
]

cw = Inches(5.9)
ch = Inches(2.4)
cx0, cy0 = Inches(0.42), Inches(1.75)

for i, (col, tag, title, body) in enumerate(cards):
    col_i = i % 2
    row_i = i // 2
    bx = cx0 + col_i * (cw + Inches(0.3))
    by = cy0 + row_i * (ch + Inches(0.22))

    card = add_rect(s4, bx, by, cw, ch, RGBColor(0x06, 0x08, 0x10))
    card.line.color.rgb = RGBColor(0x18, 0x20, 0x35)
    card.line.width = Pt(0.5)

    # Left accent bar
    add_rect(s4, bx, by, Inches(0.05), ch, col)

    inner_x = bx + Inches(0.2)
    add_text(s4, tag,   inner_x, by + Inches(0.18), cw - Inches(0.3), Inches(0.25),
             size=9, bold=True, colour=col)
    add_text(s4, title, inner_x, by + Inches(0.48), cw - Inches(0.3), Inches(0.42),
             size=14, bold=True, colour=TEXT)
    add_text(s4, body,  inner_x, by + Inches(0.96), cw - Inches(0.3), Inches(1.3),
             size=10, colour=MUTED)

add_text(s4, "04 / 06", Inches(12.1), Inches(7.0), Inches(1.1), Inches(0.35),
         size=9, colour=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE COMPANY BRAIN
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
fill_bg(s5, BG)
add_rect(s5, Inches(0), Inches(0), Inches(0.06), H, CLAY)

add_text(s5, "INSTITUTIONAL MEMORY", Inches(0.5), Inches(0.5), Inches(10),
         Inches(0.3), size=10, bold=True, colour=CLAY)
add_text(s5, "The graph that remembers so your team doesn't have to.",
         Inches(0.5), Inches(0.9), Inches(12.5), Inches(0.65),
         size=28, bold=True, colour=TEXT)

rows5 = [
    (CLAY,   "MEMORY",  "Everything compounds, nothing decays",
     "Every meeting is auto-ingested. Decisions, people, technologies, and relationships accumulate "
     "weight across sessions. The graph reflects what your company actually spends its thinking on."),
    (CYAN,   "DEBRIEF", "Cross-meeting contradiction detection",
     "After every session, a debrief agent queries the full historical graph, surfaces decisions that "
     "contradict past ones, and emails open action items before anyone has to chase them."),
    (PURPLE, "QUERY",   "Ask anything in plain language, get cited answers",
     "Natural-language questions answered from your entire meeting history. Responses cite the exact "
     "chunk they came from. A 3D force-directed graph makes the knowledge visually navigable."),
]

ry = Inches(1.75)
for col, badge, title, body in rows5:
    # Badge pill
    pill = add_rect(s5, Inches(0.5), ry + Inches(0.05), Inches(1.1), Inches(0.32), BG2)
    pill.line.color.rgb = col
    pill.line.width = Pt(0.6)
    add_text(s5, badge, Inches(0.52), ry + Inches(0.05), Inches(1.08), Inches(0.32),
             size=8, bold=True, colour=col, align=PP_ALIGN.CENTER)

    add_text(s5, title, Inches(1.75), ry,              Inches(11.0), Inches(0.38),
             size=15, bold=True, colour=TEXT)
    add_text(s5, body,  Inches(1.75), ry + Inches(0.4), Inches(11.0), Inches(0.65),
             size=11, colour=MUTED)

    # Divider
    if body != rows5[-1][3]:
        add_rect(s5, Inches(0.5), ry + Inches(1.22), Inches(12.3), Inches(0.01),
                 RGBColor(0x1A, 0x22, 0x38))
    ry += Inches(1.4)

add_text(s5, "05 / 06", Inches(12.1), Inches(7.0), Inches(1.1), Inches(0.35),
         size=9, colour=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — STACK / CTA
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
fill_bg(s6, BG2)
add_rect(s6, Inches(0), Inches(0), Inches(0.06), H, CLAY)

add_text(s6, "BUILT ON", Inches(0.5), Inches(0.55), Inches(8),
         Inches(0.3), size=10, bold=True, colour=CLAY)

add_text(s6, "Zero hardware.",    Inches(0.5), Inches(1.0),  Inches(9), Inches(0.75),
         size=48, bold=True, colour=TEXT)
add_text(s6, "Zero enrollment.",  Inches(0.5), Inches(1.75), Inches(9), Inches(0.75),
         size=48, bold=True, colour=TEXT)
add_text(s6, "Every meeting.",    Inches(0.5), Inches(2.5),  Inches(9), Inches(0.75),
         size=48, bold=True, colour=CLAY)

# Chip row
chips_lit  = ["MediaPipe", "Groq Whisper", "Neo4j", "Claude"]
chips_dim  = ["Chrome MV3", "FastAPI", "Three.js", "React", "GraphRAG", "all-MiniLM-L6-v2"]

cx = Inches(0.5)
cy = Inches(3.6)
chip_h = Inches(0.35)

for label in chips_lit + chips_dim:
    lit = label in chips_lit
    # estimate width
    cw2 = Inches(max(0.9, len(label) * 0.085 + 0.3))
    chip = add_rect(s6, cx, cy, cw2, chip_h,
                    RGBColor(0x1C, 0x0A, 0x08) if lit else BG)
    chip.line.color.rgb = CLAY if lit else RGBColor(0x1A, 0x22, 0x38)
    chip.line.width = Pt(0.6)
    add_text(s6, label, cx + Inches(0.08), cy + Inches(0.04), cw2 - Inches(0.1), chip_h,
             size=9, bold=lit, colour=CLAY if lit else MUTED, align=PP_ALIGN.CENTER)
    cx += cw2 + Inches(0.12)
    if cx > Inches(12.5):   # wrap to next row
        cx = Inches(0.5)
        cy += chip_h + Inches(0.1)

# Tagline callout
tl_y = Inches(5.3)
add_rect(s6, Inches(0.5), tl_y, Inches(0.05), Inches(0.75), CLAY)
add_text(s6, "Your meetings. Your knowledge.", Inches(0.7), tl_y,
         Inches(11.0), Inches(0.4), size=16, bold=True, colour=TEXT)
add_text(s6, "Never lost again.", Inches(0.7), tl_y + Inches(0.4),
         Inches(11.0), Inches(0.35), size=16, colour=MUTED)

add_text(s6, "06 / 06", Inches(12.1), Inches(7.0), Inches(1.1), Inches(0.35),
         size=9, colour=MUTED, align=PP_ALIGN.RIGHT)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "recognize-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
